import requests
import datetime
import os
import re
import urllib3
from datetime import datetime
from urllib3.exceptions import InsecureRequestWarning
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION
from bs4 import BeautifulSoup
from PIL import Image
from io import BytesIO
from pptx.enum.text import MSO_AUTO_SIZE
import unicodedata
import html2text
from utils.caching import BASE_URL as KEEPITSECURE_URL
from utils.manage_keys import get_credential
import flet as ft
from utils import state
from logs.logger import logger

# Preparing requests
urllib3.disable_warnings(InsecureRequestWarning)

ORGANISATIONS = {}
OUTPUT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__),  "..", "Reports"))
os.makedirs(OUTPUT_DIR, exist_ok=True)
ATTACHMENTS_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "attachments"))
x_api_key = get_credential("DashboardAPIKey")

headers = {
    "Content-Type":"application/json",
    "Accept": "application/json",
    "x-api-key": x_api_key
}

proxies = {
   "http": "127.0.0.1:8080",
   "https": "127.0.0.1:8080",
}

# Load the template PowerPoint file
template_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "templates", "presentation", "template.pptx"))
# presentation = Presentation(template_path)

# Define a dictionary of replacements
replacements = {
    "[TARGET_NAME]": "",
    "[ASSET_ID]": "",
    "[SERVICE_TYPE]": "",
    "[TIMESTAMP]": f"Generated on {datetime.now().strftime('%d-%m-%Y')}",
    "[OPCO]": "",
    "[DURATION]": "",
    "[START_DATE]": "DD-MM-YYYY",
    "[END_DATE]": "DD-MM-YYYY",
    "[TOTAL_VULNS]": "",
    "[VULNS_SUMMARY]": "",
    "[VULN_TITLE]": "",
    "[VULN_SEVERITY]": "",
    "[VULN_STATUS]": "",
    "[VULN_DESCRIPTION]": "",
    "[VULN_IMPACT]": "",
    "[VULN_RECOMMENDATION]": "",
    "[VULN_DUE]": "DD-MM-YYYY",
    "[REQUEST_ID]": "",
    "[MANAGEMENT_SUMMARY]": "",
    # Add more placeholders and their replacements as needed
}

# Service types
services = {
    "7ea682fb-e03e-43c8-a476-877d088b6cdb" : "Black Box",
    "ff52df96-1934-4cc0-9090-804dc9feeae9" : "Grey Box",
    "e14f6754-8b1f-41e7-88dd-cbb4fda7a3f2" : "White Box",
    "43b55a12-bfb3-4e2d-a212-0db9d73358e4" : "Adversary Simulation",
}

group_of_tests = {
  "#": ("Group", "E.g"),
  "1": ("Information Gathering", "Fingerprint Web Server, Fingerprint Web Application, Map Application Architecture"),
  "2": ("Configuration and Deploy Management Testing",
        "Application Platform Configuration, HSTS, HTTP Methods, Backup and Unreferenced Files"),
  "3": ("Identity Management Testing", "User Registration Process, Account Enumeration"),
  "4": ("Authentication Testing", "Default Credentials, Weak Lock Out Mechanism, Reset Password, Weak Password Policy"),
  "5": ("Authorization Testing",
        "Directory Traversal/File Inclusion, Insecure Direct Object References, Privilege Escalation"),
  "6": ("Session Management Testing",
        "Session Fixation, Logout Functionality, Session Timeout, Bypassing Session Management Schema"),
  "7": ("Data Validation Testing", "SQL Injection, Code Injection, Reflected XSS, Stored XSS, HTTP Parameter Pollution"),
  "8": ("Error Handling", "Analysis of Error Codes, Analysis of Stack Traces"),
  "9": ("Cryptography", "Weak SSL/TLS Ciphers, Sensitive information sent via unencrypted channels"),
  "10": ("Business Logic Testing",
         "Ability to Forge Requests, Business Logic Data Validation, Upload of Malicious Files"),
  "11": ("Client Side Testing", "JavaScript Execution, HTML Injection, Client Side URL Redirect, Clickjacking")
}

## Colors
gColors = {
  'darkred'  : RGBColor(168,  0,  0),
  'red'      : RGBColor(231, 69, 54),
  'orange'   : RGBColor(255,181, 17),
  'green'    : RGBColor(  0,204,102),
  'yellow'   : RGBColor(255,253,173),
  'white'    : RGBColor(255,255,255),
  'black'    : RGBColor(  0,  0,  0),
  'blue'     : RGBColor( 33,117,217),
}

gColorsRisk = {
  'Critical' : gColors['darkred'],
  'High'     : gColors['red'],
  'Medium'   : gColors['orange'],
  'Low'      : gColors['yellow'],
  'Info'     : gColors['blue']
}


# Function to add a slide
def _add_slide(presentation, slide_layout_name_to_index, layout_name):
    return presentation.slides.add_slide(
        presentation.slide_layouts[slide_layout_name_to_index[layout_name]]
    )


def retrieve_test_info(testUUID):
    body = {}
    body["uuid"] = testUUID
    session = requests.Session()
    resp = session.post(f"{KEEPITSECURE_URL}/api/v3/tests", json=body, headers=headers, verify=False)
    return resp.json()["items"][0]


def retrieve_test_service(page, testUUID):
    body = {}
    body["tests"] = testUUID
    session = requests.Session()
    resp = session.post(f"{KEEPITSECURE_URL}/api/v3/tags", json=body, headers=headers, verify=False)

    for item in resp.json()["items"]:
        try:
            if services[item["uuid"]]:
                return services[item["uuid"]]
        except Exception as e:
            logger.exception(f"Failed to retrieve service: {e}")
            pass

    logger.info("No service type found.")
    page.snack_bar.content = ft.Row(
        [
            ft.Icon(name=ft.Icons.WARNING_OUTLINED, color=ft.Colors.BLACK87),
            ft.Text("No service type found.", color=ft.Colors.BLACK87)
        ]
    )
    page.snack_bar.bgcolor = ft.Colors.ORANGE_400
    page.snack_bar.open = True
    return ""


# Getting asset data
def retrieve_asset_info(testUUID):
    body = {}
    body["tests"] = testUUID
    session = requests.Session()
    resp = session.post(f"{KEEPITSECURE_URL}/api/v3/assets", json=body, headers=headers, verify=False)

    return resp.json()["items"]


def retrieve_onetrust_id(page, assetUUID):
    body = {}
    body["assets"] = assetUUID
    session = requests.Session()
    resp = session.post(f"{KEEPITSECURE_URL}/api/v3/fields", json=body, headers=headers, verify=False)

    for item in resp.json()["items"]:
        if item["custom_field"]["name"] == "OneTrust ID":
            if isinstance(item["value"], str):
                return item["value"]
            else:
                logger.info("Asset is missing OneTrust ID value")
                page.snack_bar.content = ft.Row(
                    [
                        ft.Icon(name=ft.Icons.WARNING_OUTLINED, color=ft.Colors.BLACK87),
                        ft.Text("Asset is missing OneTrust ID value", color=ft.Colors.BLACK87)
                    ]
                )
                page.snack_bar.bgcolor = ft.Colors.ORANGE_400
                page.snack_bar.open = True
                return ""

    return ""


# Getting vulnerabilities data
def retrieve_vulns_by_test(testUUID):
    body = {}
    body["tests"] = testUUID
    session = requests.Session()
    resp = session.post(f"{KEEPITSECURE_URL}/api/v3/vulnerabilities", json=body, headers=headers, verify=False)

    return resp.json()["items"]


def _add_image(presentation, slide_layout_name_to_index, image_url):
    slide = _add_slide(presentation, slide_layout_name_to_index, "title only")
    slide.placeholders[0].text = "evidence"

    # Calculate the image size of the image
    image_content = open(f'{image_url}', 'rb').read()
    im = Image.open(BytesIO(image_content))
    width, height = im.size

    pic_width = int((presentation.slide_width - 100) * 0.8)
    pic_height = int((presentation.slide_height - 100) * 0.8)

    if width >= height:
        picture = slide.shapes.add_picture(f'{image_url}', Inches(1.5), Inches(1.5), width=pic_width)
        picture.crop_left = 0
        picture.crop_right = 0
    else:
        picture = slide.shapes.add_picture(f'{image_url}', Inches(1.5), Inches(1.5), height=pic_height)
        picture.crop_top = 0
        picture.crop_bottom = 0
    ratio = min(pic_width / float(picture.width), pic_height / float(picture.height))
    picture.height = int(picture.height * ratio)
    picture.width = int(picture.width * ratio)
    picture.left = int((presentation.slide_width - picture.width) / 2)


def retrieve_vuln_attachments(presentation, slide_layout_name_to_index, vulnUUID):
    session = requests.Session()

    # 1) fetch the vulnerability, so we get back the attachments metadata
    resp = session.post(
        f"{KEEPITSECURE_URL}/api/v3/vulnerabilities",
        headers=headers,
        json={"uuid": [vulnUUID]},
        verify=False
    )
    vulns = resp.json()["items"]
    if not vulns:
        return

    att_meta = vulns[0]["attachments"]

    # 2) now actually ask for the attachment records
    resp2 = session.post(
        f"{KEEPITSECURE_URL}{att_meta['url']}",
        headers=headers,
        json=att_meta["body"],
        verify=False
    )
    attachment_list = resp2.json().get("items", [])
    logger.debug(f"FOUND {len(attachment_list)} ATTACHMENTS")

    # 3) iterate the real attachments
    for attachment in attachment_list:
        uuid = attachment["uuid"]
        out_path = os.path.join(ATTACHMENTS_DIR, f"{uuid}.png")

        if not os.path.exists(out_path):
            img_resp = session.get(
                f"{KEEPITSECURE_URL}/api/v3/attachments/{uuid}",
                headers=headers,
                verify=False
            )
            with open(out_path, "wb") as f:
                f.write(img_resp.content)

        if is_image_file(out_path):
            _add_image(presentation, slide_layout_name_to_index, out_path)


def is_image_file(filename):
    with open(filename, 'rb') as f:
        header = f.read(4)  # Read the first 4 bytes (header)
        if header.startswith(b'\x89\x50\x4E\x47'):  # PNG magic number
            return True
        elif header.startswith(b'\xFF\xD8'):  # JPEG/JPG magic number
            return True
        else:
            return False


def vulns_summary(presentation, vulns):
    critical = high = medium = low = info = 0

    for vuln in vulns:
        if vuln["severity"] == "Critical":
            critical = critical + 1
        if vuln["severity"] == "High":
            high = high + 1
        if vuln["severity"] == "Medium":
            medium = medium + 1
        if vuln["severity"] == "Low":
            low = low + 1
        if vuln["severity"] == "Info":
            info = info + 1

    chart(presentation, 5, critical, high, medium, low, info)
    total = {"Critical": critical, "High" : high, "Medium" : medium, "Low" : low, "Info" : info}
    total = {x:y for x,y in total.items() if y != 0} # Removing empty severity

    message = ", including "

    if len(total) > 0:
        for index, (key, value) in enumerate(total.items()):
            if len(total) == 1:
                message = message + f'{total[key]} {key} Severity'
            else:
                if index == len(total) - 1:
                    message = message + f'and {total[key]} {key} Severity'
                else:
                    message = message + f'{total[key]} {key} Severity, '

        message = message + ". The total number of vulnerabilities in each category and applications has been summarized as follows:"
    else:
        message = "."

    return message


# Function to replace placeholders with text
def replace_text(slide, replacements):
    logger.info("Replacing text on slide:")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for old_text, new_text in replacements.items():
                        if old_text in run.text:
                            logger.info(f"  Found placeholder: {old_text} → {new_text}")
                            run.text = run.text.replace(old_text, new_text)


def chart(presentation, slide_index, critical, high, medium, low, info):
    ## Define chart data # Creating object of chart
    chart_data = CategoryChartData()
    ## Adding categories to chart
    chart_data.categories = ['']
    ## Adding series
    chart_data.add_series('Critical', (int(critical),))
    chart_data.add_series('High',     (int(high),))
    chart_data.add_series('Medium',   (int(medium),))
    chart_data.add_series('Low',      (int(low),))
    chart_data.add_series('Info',     (int(info),))
    x, y, cx, cy = Inches(2.7), Inches(2.5), Inches(8), Inches(4.5)
    slide = presentation.slides[slide_index]
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data )
    chart = graphic_frame.chart
    colors = {
        0: gColorsRisk['Critical'],
        1: gColorsRisk['High'],
        2: gColorsRisk['Medium'],
        3: gColorsRisk['Low'],
        4: gColorsRisk['Info'],
    }
    for idx,serie in enumerate(chart.plots[0].series):
      fill = serie.format.fill
      fill.solid()
      fill.fore_color.rgb = colors[idx]
    ## Data label
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(16)
    data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    #
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(13)
    chart.legend.include_in_layout = False
    # Overlap/Gap Width
    plot = chart.plots[0]
    plot.gap_width = 0
    plot.overlap = -5

    # Customize the font size of the vertical axis values (numbers)
    vertical_axis = chart.value_axis  # Get the vertical axis
    vertical_axis.has_major_gridlines = False  # Optional: Hide major gridlines

    # Access the font properties of the vertical axis
    vertical_axis.tick_labels.font.size = Pt(10)

    # Set the vertical axis major unit to 1 (display only integers)
    vertical_axis.major_unit = 1

    # Customize the gap width between bars (adjust as needed)
    plot = chart.plots[0]
    plot.gap_width = 200  # Increase this value to widen the gap between bars (percentage)


def slugify(value, allow_unicode=False):
    """
    Taken from https://github.com/django/django/blob/master/django/utils/text.py
    Convert to ASCII if 'allow_unicode' is False. Convert spaces or repeated
    dashes to single dashes. Remove characters that aren't alphanumerics,
    underscores, or hyphens. Convert to lowercase. Also strip leading and
    trailing whitespace, dashes, and underscores.
    """
    value = str(value)
    if allow_unicode:
        value = unicodedata.normalize('NFKC', value)
    else:
        value = unicodedata.normalize('NFKD', value).encode('ascii', 'ignore').decode('ascii')
    value = re.sub(r'[^\w\s-]', '_', value)

    return value


def main(page, e):
    main_headers = {
        "Content-Type": "application/x-www-form-urlencoded",
        "Accept": "application/json",
        "x-api-key": x_api_key
    }

    # Create a session object
    session = requests.Session()

    if not os.path.exists(OUTPUT_DIR):
        os.mkdir(OUTPUT_DIR)

    if not os.path.exists(ATTACHMENTS_DIR):
        os.mkdir(ATTACHMENTS_DIR)

    presentation = Presentation(template_path)
    slide_layout_name_to_index = {
        layout.name: i for i, layout in enumerate(presentation.slide_layouts)
    }

    testUUID = page.app_state.presentation_test_uuid_text_field.value

    if not testUUID:
        page.snack_bar.content = ft.Row(
            [
                ft.Icon(name=ft.Icons.WARNING_OUTLINED, color=ft.Colors.BLACK87),
                ft.Text("Test UUID not Found!", color=ft.Colors.BLACK87)
            ]
        )
        page.snack_bar.bgcolor = ft.Colors.ORANGE_400

    test = retrieve_test_info(testUUID)
    test_service = retrieve_test_service(page, testUUID)
    asset = retrieve_asset_info(testUUID)[0]
    onetrust_id = retrieve_onetrust_id(page, asset["uuid"])

    test_start = datetime.strptime(test["started_at"], "%Y-%m-%d %H:%M:%S")

    test_end = None
    test_duration = None
    if isinstance(test, dict):
        raw = test.get("ended_at")
        if isinstance(raw, str):
            test_end = datetime.strptime(raw, "%Y-%m-%d %H:%M:%S")
            test_duration = test_end - test_start

    if test_end is None:
        page.snack_bar.content = ft.Row(
            [
                ft.Icon(name=ft.Icons.WARNING_OUTLINED, color=ft.Colors.BLACK87),
                ft.Text("The test is not Ended yet. End the test before generate the Presentation!", color=ft.Colors.BLACK87)
            ]
        )
        page.snack_bar.bgcolor = ft.Colors.ORANGE_400
        page.snack_bar.open = True
        page.update()

    if test_duration.days <= 0:
        replacements["[DURATION]"] = "1 day"
    else:
        replacements["[DURATION]"] = str(test_duration.days) + " days"

    replacements["[OPCO]"] = test["organisation"]["name"]
    replacements["[SERVICE_TYPE]"] = test_service
    replacements["[TARGET_NAME]"] = asset["name"]
    replacements["[ASSET_ID]"] = onetrust_id
    replacements["[START_DATE]"] = test_start.strftime("%d/%m/%Y")
    replacements["[END_DATE]"] = test_end.strftime("%d/%m/%Y")

    # To remove html tags
    config = html2text.HTML2Text()
    config.body_width = 0  # disable line wrapping
    _txt = config.handle(test['details'])
    # Extract ServiceNow Request ID using regex
    service_now_request_id = re.search(r"ServiceNow Request ID: ([^\s\t\n\r]+)", _txt)
    if service_now_request_id:
        replacements["[REQUEST_ID]"] = service_now_request_id.group(1)
    else:
        logger.info("ServiceNow Request ID not found.")

    # Extract content after "[MANAGEMENT SUMMARY]" using regex
    management_summary_content = re.search(r'[\[{]MANAGEMENT SUMMARY[\]}].*?[\s\t\r\n]+(.*)', _txt, re.DOTALL|re.I)

    # Check if the regex match was successful
    if management_summary_content:
        replacements["[MANAGEMENT_SUMMARY]"] = "\n".join(management_summary_content.group(1).splitlines())
    else:
        logger.info("Management summary content not found.")

    vulns = retrieve_vulns_by_test(testUUID)

    replacements["[TOTAL_VULNS]"] = str(len(vulns))
    replacements["[VULNS_SUMMARY]"] = vulns_summary(presentation, vulns)

    # Replace placeholders in each slide
    for slide in presentation.slides:
        replace_text(slide, replacements)

    #################################
    # Adding LIST OF FINDINGS slide #
    #################################
    slide = _add_slide(presentation, slide_layout_name_to_index, "title only")
    slide.placeholders[0].text = "list of findings"

    # Define the number of rows and columns for the table
    num_rows = len(vulns) + 1
    num_cols = 3

    # Define the width and height of the table (in inches)
    table_width = Inches(11)
    table_height = Inches(1)

    # Define the position of the table on the slide (left, top) in inches
    left = Inches(0.8)
    top = Inches(1.7)

    # Insert a table into the slide
    shape = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height)
    table = shape.table
    table.columns[0].width = Inches(7)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(2)

    tbl =  shape._element.graphic.graphicData.tbl
    style_id = "{D3789944-8FE3-4B66-81FE-070CC84C6259}"
    tbl[0][-1].text = style_id

    # Populate the cells of the table with content
    cell = table.cell(0, 0)
    cell.text = "finding"
    cell = table.cell(0, 1)
    cell.text = "severity"
    cell = table.cell(0, 2)
    cell.text = "status"

    row = 1

    # Sort all vulnerabilities by descending
    severity_order = [
        "info",
        "low",
        "medium",
        "high",
        "critical"
    ]
    severity_order = severity_order[::-1]
    sorted_vulns = list()
    for severity in severity_order:
        for vuln in vulns:
            if vuln['severity'].lower() == severity:
                sorted_vulns.append(vuln)

    vulns = sorted_vulns

    for vuln in vulns:
        cell = table.cell(row, 0)
        cell.text = f"{vuln['description']}"

        cell = table.cell(row, 1)
        cell.text = f"{vuln['severity']}"

        cell = table.cell(row, 2)
        if vuln['state'] == "Unpublished" or vuln['state'] == "Ready To Publish" or vuln['state'] == "New":
            cell.text = "Open"
        else:
            cell.text = f"{vuln['state']}"

        row = row + 1

    ############################
    # Adding VULN DETAIL slide #
    ############################
    for vuln in vulns:
        slide = _add_slide(presentation, slide_layout_name_to_index, "detailed vuln")
        a, b, c, d, e, f, g, h = slide.placeholders
        a.text_frame.word_wrap = True
        a.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        b.text_frame.word_wrap = True
        b.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        d.text_frame.word_wrap = True
        d.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        f.text_frame.word_wrap = True
        f.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        h.text_frame.word_wrap = True
        h.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        a.text = vuln["description"]
        b.text = vuln["severity"]
        g.text = "description"
        c.text = "impact"
        e.text = "recommendation"
        fill = b.fill
        fill.solid()
        if b.text == "Critical":
            fill.fore_color.rgb = gColorsRisk["Critical"]
        elif b.text == "High":
            fill.fore_color.rgb = gColorsRisk["High"]
        elif b.text == "Medium":
            fill.fore_color.rgb = gColorsRisk["Medium"]
        elif b.text == "Low":
            fill.fore_color.rgb = gColorsRisk["Low"]
            font = b.text_frame.paragraphs[0].runs[0].font
            font.color.rgb = RGBColor(0, 0, 0)
        elif b.text == "Info":
            fill.fore_color.rgb = gColorsRisk["Info"]

        # Parse the HTML content
        soup = BeautifulSoup(vuln['details'], 'html.parser')

        # Extract the content of each section
        sections = {}
        current_section = None

        for tag in soup.find_all(['h1', 'p', 'li']):
            if tag.name == 'h1':
                current_section = tag.text.strip()
                sections[current_section] = ""
            elif (tag.name == 'p' or tag.name == 'li') and current_section:
                sections[current_section] += tag.text.strip() + "\n"

        # Print or store the content of each section
        h.text = "\n".join(sections["Description"].splitlines())
        d.text = "\n".join(sections["Impact"].splitlines())
        f.text = "\n".join(sections["Recommendation"].splitlines())

        if int(vuln["attachments"]["total"]) > 0:
            retrieve_vuln_attachments(presentation, slide_layout_name_to_index, vuln["uuid"])

    ##################
    # Adding divider #
    ##################
    slide = _add_slide(presentation, slide_layout_name_to_index, "divider dark blue")
    slide.placeholders[0].text = "what's next"

    ##################################
    # Adding FINDINGS DUE DATE slide #
    ##################################
    slide = _add_slide(presentation, slide_layout_name_to_index, "title only")
    slide.placeholders[0].text = "findings due date overview"

    # Define the number of rows and columns for the table
    num_rows = len(vulns) + 1
    num_cols = 3

    # Define the width and height of the table (in inches)
    table_width = Inches(11)
    table_height = Inches(1)

    # Define the position of the table on the slide (left, top) in inches
    left = Inches(0.8)
    top = Inches(1.7)

    # Insert a table into the slide
    shape = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height)
    table = shape.table
    table.columns[0].width = Inches(7)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(2)

    tbl = shape._element.graphic.graphicData.tbl
    style_id = "{D3789944-8FE3-4B66-81FE-070CC84C6259}"
    tbl[0][-1].text = style_id

    # Populate the cells of the table with content
    cell = table.cell(0, 0)
    cell.text = "finding"
    cell = table.cell(0, 1)
    cell.text = "severity"
    cell = table.cell(0, 2)
    cell.text = "due date"

    row = 1

    if state.selected_date_str:
        now = datetime.strptime(state.selected_date_str, "%d-%m-%Y")
    else:
        now = datetime.now()

    for vuln in vulns:
        cell = table.cell(row, 0)
        cell.text = vuln['description']

        cell = table.cell(row, 1)
        cell.text = vuln['severity']

        cell = table.cell(row, 2)
        if vuln['severity'] == "Critical":
          #due_date = now + timedelta(days=14)
          #cell.text = due_date.strftime('%d-%m-%Y')
          cell.text = "14 days"
        elif vuln['severity'] == "High":
          #due_date = now + timedelta(days=30)
          #cell.text = due_date.strftime('%d-%m-%Y')
          cell.text = "30 days"
        elif vuln['severity'] == "Medium":
          #due_date = now + timedelta(days=45)
          #cell.text = due_date.strftime('%d-%m-%Y')
          cell.text = "45 days"
        elif vuln['severity'] == "Low":
          #due_date = now + timedelta(days=60)
          #cell.text = due_date.strftime('%d-%m-%Y')
          cell.text = "60 days"
        elif vuln['severity'] == "Info":
          #due_date = now + timedelta(days=270)
          #cell.text = due_date.strftime('%d-%m-%Y')
          cell.text = "9 months"

        row = row + 1

    ####################
    # Adding questions #
    ####################
    slide = _add_slide(presentation, slide_layout_name_to_index, "questions")
    a, b = slide.placeholders
    a.text = "q & a"
    b.text = "any questions?"

    ##################
    # Adding divider #
    ##################
    slide = _add_slide(presentation, slide_layout_name_to_index, "divider dark blue")
    slide.placeholders[0].text = "appendix"

    #########################
    # Adding APPENDIX slide #
    #########################
    slide = _add_slide(presentation, slide_layout_name_to_index, "title subtitle")
    a, b = slide.placeholders
    a.text = "group of tests performed"
    b.text = "example of tests performed during the pentest campaign"

    # Define the number of rows and columns for the table
    num_rows = len(group_of_tests)
    num_cols = 3

    # Define the width and height of the table (in inches)
    table_width = Inches(11.5)
    table_height = Inches(1)

    # Define the position of the table on the slide (left, top) in inches
    left = Inches(0.8)
    top = Inches(2)

    # Insert a table into the slide
    shape = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height)
    table = shape.table
    table.columns[0].width = Inches(0.5)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(8)

    tbl =  shape._element.graphic.graphicData.tbl
    style_id = "{D3789944-8FE3-4B66-81FE-070CC84C6259}"
    tbl[0][-1].text = style_id

    # Populate the cells of the table with content
    for index, key in enumerate(group_of_tests):
        cell = table.cell(index, 0)
        cell.text = key
        cell = table.cell(index, 1)
        cell.text = group_of_tests[key][0]
        cell = table.cell(index, 2)
        cell.text = group_of_tests[key][1]

    #############
    # Last page #
    #############
    slide = _add_slide(presentation, slide_layout_name_to_index, "last page")

    # Save the updated presentation

    test_pp = test['id']
    ass_pp = asset['name']
    date_pp = now.strftime('%Y-%m')

    output_path = os.path.join(
        OUTPUT_DIR,
        f"{test_pp}-Restitution_Meeting-{slugify(ass_pp)}-{date_pp}.pptx"
    )

    if os.path.exists(output_path):
        os.remove(output_path)

    presentation.save(output_path)
    page.snack_bar.content = ft.Row(
        [
            ft.Icon(name=ft.Icons.CHECK_OUTLINED, color=ft.Colors.BLACK87),
            ft.Text(f"Presentation generated and stored at: '{output_path}'",
                    color=ft.Colors.BLACK87)
        ]
    )
    page.snack_bar.bgcolor = ft.Colors.GREEN_400
    page.snack_bar.open = True
    page.update()
